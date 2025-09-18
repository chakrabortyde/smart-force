from typing import List
from langchain_core.documents import Document
import os

# Simple, dependency-light loaders
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx"}

def load_file(path: str) -> List[Document]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return load_pdf(path)
    if ext == ".docx":
        return load_docx(path)
    if ext == ".pptx":
        return load_pptx(path)
    return []

def load_pdf(path: str) -> List[Document]:
    reader = PdfReader(path)
    texts = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            texts.append((i + 1, text))
    return [Document(page_content=t, metadata={"source": path, "page": p}) for p, t in texts]

def load_docx(path: str) -> List[Document]:
    doc = DocxDocument(path)
    text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    if not text.strip():
        return []
    return [Document(page_content=text, metadata={"source": path})]

def load_pptx(path: str) -> List[Document]:
    prs = Presentation(path)
    slides_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    chunks.append(t)
        if chunks:
            slides_text.append((idx, "\n".join(chunks)))
    return [Document(page_content=text, metadata={"source": path, "slide": slide_no})
            for slide_no, text in slides_text]
